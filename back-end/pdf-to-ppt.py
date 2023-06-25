import os
from pdf2image import convert_from_path
from pptx import Presentation
#pip install pdf2image python-pptx TO INSTALL

def convert_pdf_to_ppt(pdf_path):
    # Convert PDF pages to images
    images = convert_from_path(pdf_path)

    # Create a new PowerPoint presentation
    ppt = Presentation()

    for i, image in enumerate(images):
        # Determine slide cutoffs based on image dimensions or custom logic
        slide_cutoff = determine_slide_cutoff(image)

        # Split the image into two parts based on the slide cutoff
        image_top = image.crop((0, 0, image.width, slide_cutoff))
        image_bottom = image.crop((0, slide_cutoff, image.width, image.height))

        # Create a new slide in the PowerPoint presentation
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # Change the layout index if needed

        # Add the top part of the image to the slide
        slide.shapes.add_picture(image_top, 0, 0, width=ppt.slide_width, height=slide_cutoff)

        # Add the bottom part of the image to the next slide (if applicable)
        if i < len(images) - 1:
            next_slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # Change the layout index if needed
            next_slide.shapes.add_picture(image_bottom, 0, 0, width=ppt.slide_width, height=image_bottom.height)

    # Save the PowerPoint presentation
    ppt_path = os.path.splitext(pdf_path)[0] + '.pptx'
    ppt.save(ppt_path)

    return ppt_path

def determine_slide_cutoff(image):
    # Implement your logic to determine the appropriate cutoff for slides
    # You can use image dimensions, text detection, or other criteria

    # Example: Split the image into two equal parts
    slide_cutoff = image.height // 2

    return slide_cutoff

# Example usage
pdf_file = '/path/to/input.pdf'
ppt_file = convert_pdf_to_ppt(pdf_file)